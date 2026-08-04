"""
Script tự động khởi tạo file PowerPoint (.pptx) thuyết trình chuyên nghiệp 12 Slides.
Dự án: Shopee AI Help Center Assistant (RAG Pipeline)
Nhóm: Tên Nhóm Là Sao
Diễn giả: Nguyễn Hùng Mạnh
"""

import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

if sys.platform == "win32":
    try:
        sys.stdout.reconfigure(encoding="utf-8")
    except Exception:
        pass

# Color Palette Constants (Shopee Theme)
COLOR_ORANGE = RGBColor(238, 77, 45)      # Shopee Primary #EE4D2D
COLOR_DARK = RGBColor(15, 23, 42)         # Slate 900 #0F172A
COLOR_NAVY = RGBColor(30, 41, 59)         # Slate 800 #1E293B
COLOR_GRAY = RGBColor(241, 245, 249)      # Slate 100 #F1F5F9
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_MUTED = RGBColor(100, 116, 139)     # Slate 500 #64748B
COLOR_CORAL = RGBColor(255, 107, 0)       # Accent Orange #FF6B00

def create_deck():
    prs = Presentation()
    # Set 16:9 Widescreen Layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    def add_header(slide, title_text, category_text="SHOPEE RAG AI DỰ ÁN K4"):
        # Header Background Banner
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
        header.fill.solid()
        header.fill.fore_color.rgb = COLOR_DARK
        header.line.fill.background()

        # Category Badge
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.15), Inches(11), Inches(0.3))
        tf_cat = cat_box.text_frame
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = COLOR_ORANGE

        # Main Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(11), Inches(0.55))
        tf_title = title_box.text_frame
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(22)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_WHITE

    # -------------------------------------------------------------------------
    # SLIDE 1: Title Slide (Dark Theme)
    # -------------------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = COLOR_DARK
    bg1.line.fill.background()

    # Accent Stripe
    stripe = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(0.15), Inches(4.5))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = COLOR_ORANGE
    stripe.line.fill.background()

    t_box = slide1.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(11.3), Inches(4.5))
    tf1 = t_box.text_frame
    tf1.word_wrap = True

    p1 = tf1.paragraphs[0]
    p1.text = "SHOPEE AI HELP CENTER ASSISTANT"
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_WHITE

    p2 = tf1.add_paragraph()
    p2.text = "Hệ thống RAG Tìm kiếm Hybrid Kết hợp Sinh câu trả lời có Trích dẫn Nguồn"
    p2.font.size = Pt(22)
    p2.font.color.rgb = COLOR_ORANGE
    p2.space_before = Pt(15)

    p3 = tf1.add_paragraph()
    p3.text = "Nhóm Thực Hiện: Tên Nhóm Là Sao | Diễn giả: Nguyễn Hùng Mạnh"
    p3.font.size = Pt(16)
    p3.font.color.rgb = COLOR_MUTED
    p3.space_before = Pt(30)

    # -------------------------------------------------------------------------
    # SLIDE 2: Problem Statement & Objectives
    # -------------------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    add_header(slide2, "BỐI CẢNH DỰ ÁN & THÁCH THỨC ĐẶT RA")

    # Card 1: Problem
    card1 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.2))
    card1.fill.solid()
    card1.fill.fore_color.rgb = COLOR_GRAY
    card1.line.color.rgb = COLOR_MUTED

    tf = card1.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "⚠️ Thách Thức của CSKH TMĐT"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_DARK

    bullets1 = [
        "Hàng triệu câu hỏi trùng lặp về đổi trả, thanh toán, gian lận mỗi ngày.",
        "LLM thuần túy bị 'Bịa đặt' (Hallucination) và không có tri thức Shopee.",
        "LLM thuần túy không thể chỉ ra căn cứ điều khoản trích dẫn nguồn.",
        "Tìm kiếm từ khóa truyền thống không hiểu được ý định tự nhiên của khách."
    ]
    for b in bullets1:
        p = tf.add_paragraph()
        p.text = f"• {b}"
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_NAVY
        p.space_before = Pt(12)

    # Card 2: Solution
    card2 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(5.6), Inches(5.2))
    card2.fill.solid()
    card2.fill.fore_color.rgb = COLOR_GRAY
    card2.line.color.rgb = COLOR_ORANGE

    tf = card2.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🎯 Mục Tiêu Giải Pháp RAG"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_ORANGE

    bullets2 = [
        "Chính xác 100%: Trả lời hoàn toàn dựa trên kho tri thức Shopee chính thức.",
        "Trích dẫn minh chứng: Tự động gắn tag nguồn [Document X | Source: ...].",
        "Phản hồi siêu tốc: Retrieval trong < 0.05 giây nhờ Pre-warming RAM.",
        "Streaming thời gian thực: Hiển thị chữ gõ mượt mà từng từ như ChatGPT."
    ]
    for b in bullets2:
        p = tf.add_paragraph()
        p.text = f"• {b}"
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_NAVY
        p.space_before = Pt(12)

    # -------------------------------------------------------------------------
    # SLIDE 3: End-to-End Architecture
    # -------------------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    add_header(slide3, "KIẾN TRÚC TỔNG QUAN HỆ THỐNG RAG 10 TASKS")

    steps = [
        ("1. Data Engineering", "Task 1-3: Crawl News & PDF Legal Cleaning"),
        ("2. Vector Indexing", "Task 4: Recursive Chunking 800/100 & BAAI/bge-m3"),
        ("3. Dense Search", "Task 5: Vector Similarity + HyDE Generator"),
        ("4. Sparse Search", "Task 6: BM25Okapi Exact Keyword Match"),
        ("5. Rank Fusion", "Task 7-9: RRF Fusion (k=60) + PageIndex Fallback"),
        ("6. Generation", "Task 10: Citation Prompt + 4-Tier Fallback Chain")
    ]

    for i, (title, desc) in enumerate(steps):
        col = i % 3
        row = i // 3
        x = Inches(0.8 + col * 3.9)
        y = Inches(1.5 + row * 2.7)

        box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.6), Inches(2.3))
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_DARK if i in (2,3,4) else COLOR_GRAY
        box.line.color.rgb = COLOR_ORANGE

        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_ORANGE if i in (2,3,4) else COLOR_DARK

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = COLOR_WHITE if i in (2,3,4) else COLOR_NAVY
        p2.space_before = Pt(10)

    # -------------------------------------------------------------------------
    # SLIDE 4: Data Engineering & Cleaning (Task 1-3)
    # -------------------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    add_header(slide4, "LÀM SẠCH VÀ CHUẨN HÓA DỮ LIỆU (TASK 1 - TASK 3)")

    left_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.2))
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = COLOR_GRAY
    tf = left_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "📂 Nguồn Dữ Liệu Thực Tế"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_DARK

    items = [
        "3 Văn bản Pháp luật PDF chính thức của Shopee (Phương thức thanh toán, Chính sách bảo mật, Trả hàng hoàn tiền).",
        "5 Bài viết Tin tức hỗ trợ crawl từ help.shopee.vn (Format JSON).",
        "Tổng cộng 8 tài liệu gốc được chuyển sang Markdown (.md)."
    ]
    for it in items:
        p = tf.add_paragraph()
        p.text = f"• {it}"
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_NAVY
        p.space_before = Pt(14)

    right_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(5.6), Inches(5.2))
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = COLOR_GRAY
    tf = right_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "✨ Kỹ Thuật Cleaning Đột Phá"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_ORANGE

    items = [
        "Loại bỏ 100% ký tự điều khiển ngắt trang PDF (\\x0c).",
        "Loại bỏ các dòng Header/Footer lặp lại (Trang X | Nguon...).",
        "Tự động nối liền các câu bị ngắt dòng giữa chừng do định dạng PDF.",
        "Bảo toàn nguyên vẹn cấu trúc câu văn giúp Chunking đạt hiệu quả tối đa."
    ]
    for it in items:
        p = tf.add_paragraph()
        p.text = f"• {it}"
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_NAVY
        p.space_before = Pt(14)

    # -------------------------------------------------------------------------
    # SLIDE 5: Chunking & Metadata (Task 4)
    # -------------------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    add_header(slide5, "CẮT ĐOẠN CHUNK & METADATA ĐA CHIỀU (TASK 4)")

    # Stat 1
    s1 = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(3.6), Inches(1.8))
    s1.fill.solid()
    s1.fill.fore_color.rgb = COLOR_DARK
    tf = s1.text_frame
    p = tf.paragraphs[0]
    p.text = "141 Chunks"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_ORANGE
    p = tf.add_paragraph()
    p.text = "Tổng số Chunks trong CSDL"
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_WHITE

    # Stat 2
    s2 = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.8), Inches(1.5), Inches(3.6), Inches(1.8))
    s2.fill.solid()
    s2.fill.fore_color.rgb = COLOR_DARK
    tf = s2.text_frame
    p = tf.paragraphs[0]
    p.text = "800 / 100"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_ORANGE
    p = tf.add_paragraph()
    p.text = "CHUNK_SIZE / OVERLAP (Chars)"
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_WHITE

    # Stat 3
    s3 = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(1.5), Inches(3.6), Inches(1.8))
    s3.fill.solid()
    s3.fill.fore_color.rgb = COLOR_DARK
    tf = s3.text_frame
    p = tf.paragraphs[0]
    p.text = "1024 Dim"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_ORANGE
    p = tf.add_paragraph()
    p.text = "BAAI/bge-m3 Embedding Vector"
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_WHITE

    # Role breakdown box
    r_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.6), Inches(11.6), Inches(3.1))
    r_box.fill.solid()
    r_box.fill.fore_color.rgb = COLOR_GRAY
    tf = r_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "🏷️ Phân Loại Metadata Theo Vai Trò Khách Hàng (Customer Role)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_DARK

    roles = [
        "🛒 buyer (Người Mua): 22 chunks — Hướng dẫn thanh toán, theo dõi đơn hàng, áp mã giảm giá.",
        "🏪 seller (Người Bán): 39 chunks — Quy định xử phạt gian lận, chính sách Người Bán vi phạm.",
        "🤝 both (Cả Hai): 80 chunks — Chính sách bảo mật, chính sách đổi trả hoàn tiền chung."
    ]
    for r in roles:
        p = tf.add_paragraph()
        p.text = f"• {r}"
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_NAVY
        p.space_before = Pt(10)

    # -------------------------------------------------------------------------
    # SLIDE 6: Hybrid Search Engines (Task 5 & 6)
    # -------------------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    add_header(slide6, "ĐỘNG CƠ TÌM KIẾM DENSE VÀ SPARSE (TASK 5 & TASK 6)")

    b1 = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.2))
    b1.fill.solid()
    b1.fill.fore_color.rgb = COLOR_GRAY
    tf = b1.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🧠 Dense Search + HyDE (Task 5)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_ORANGE

    p_items = [
        "Tìm kiếm theo Tương đồng Ngữ nghĩa bằng Cosine Similarity trong ChromaDB.",
        "Công nghệ HyDE (Hypothetical Document Embeddings): Dùng LLM sinh ra tài liệu giả định trước.",
        "Dùng Vector tài liệu giả định đi quét CSDL ➔ Giúp hiểu câu hỏi tự nhiên ngay cả khi không trùng từ khóa."
    ]
    for pi in p_items:
        p = tf.add_paragraph()
        p.text = f"• {pi}"
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_NAVY
        p.space_before = Pt(12)

    b2 = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(5.6), Inches(5.2))
    b2.fill.solid()
    b2.fill.fore_color.rgb = COLOR_GRAY
    tf = b2.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🔤 Sparse Lexical Search (Task 6)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_DARK

    p_items2 = [
        "Sử dụng thuật toán BM25Okapi tính tần suất xuất hiện từ khóa chính xác.",
        "Giải quyết triệt để điểm yếu của Vector Search đối với các từ khóa hiếm hoặc mã chính sách.",
        "Đảm bảo tìm chính xác 100% các từ khóa như SPayLater, NAPAS, Apple Pay, COD."
    ]
    for pi in p_items2:
        p = tf.add_paragraph()
        p.text = f"• {pi}"
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_NAVY
        p.space_before = Pt(12)

    # -------------------------------------------------------------------------
    # SLIDE 7: RRF Rank Fusion & Fallback (Task 7-9)
    # -------------------------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    add_header(slide7, "HỢP NHẤT RRF & CHỐNG LẠC ĐỀ (TASK 7 - TASK 9)")

    c_box = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.6), Inches(2.2))
    c_box.fill.solid()
    c_box.fill.fore_color.rgb = COLOR_DARK
    tf = c_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "📐 Thuật Toán Reciprocal Rank Fusion (RRF với k = 60)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_ORANGE

    p = tf.add_paragraph()
    p.text = "RRF_Score(d) = 1 / (60 + Rank_Dense(d)) + 1 / (60 + Rank_BM25(d))"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.space_before = Pt(8)

    p = tf.add_paragraph()
    p.text = "Bỏ qua sự lệch thang điểm gốc giữa Vector và BM25, cộng điểm thưởng mượt mà theo Thứ hạng để chọn ra Top 5 Chunks được CẢ 2 BÊN cùng đánh giá tốt nhất."
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_GRAY
    p.space_before = Pt(8)

    f_box = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.0), Inches(11.6), Inches(2.7))
    f_box.fill.solid()
    f_box.fill.fore_color.rgb = COLOR_GRAY
    tf = f_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "🛡️ Cơ Chế Phòng Hộ Khẩn Cấp (Task 8: PageIndex Fallback)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_DARK

    f_bullets = [
        "Kiểm tra điểm Cosine gốc của Dense Search. Nếu SCORE_THRESHOLD < 0.3 (câu hỏi hoàn toàn lạc đề/không có trong kho tri thức).",
        "Tự động kích hoạt PageIndex Fallback để bảo vệ hệ thống không trả về kết quả rác.",
        "Giúp hệ thống đạt độ tin cậy tuyệt đối trong môi trường sản xuất."
    ]
    for fb in f_bullets:
        p = tf.add_paragraph()
        p.text = f"• {fb}"
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_NAVY
        p.space_before = Pt(8)

    # -------------------------------------------------------------------------
    # SLIDE 8: Grounded LLM Generation & Citation (Task 10)
    # -------------------------------------------------------------------------
    slide8 = prs.slides.add_slide(blank_layout)
    add_header(slide8, "SINH CÂU TRẢ LỜI CÓ TRÍCH DẪN (TASK 10)")

    # 4 Pillars
    p1_box = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.6), Inches(2.5))
    p1_box.fill.solid()
    p1_box.fill.fore_color.rgb = COLOR_GRAY
    tf = p1_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🧠 Anti 'Lost in the Middle'"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_ORANGE
    p = tf.add_paragraph()
    p.text = "Đảo thứ tự 5 Chunks thành [1, 3, 5, 4, 2] để đưa thông tin quan trọng nhất lên đầu và cuối Prompt, giúp LLM không bị bỏ sót thông tin ở giữa."
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_NAVY
    p.space_before = Pt(6)

    p2_box = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(5.6), Inches(2.5))
    p2_box.fill.solid()
    p2_box.fill.fore_color.rgb = COLOR_GRAY
    tf = p2_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🏷️ Citation Grounding & Prompting"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_DARK
    p = tf.add_paragraph()
    p.text = "Gắn thẻ [Document X | Source: ...] và đặt Temperature = 0.3 ép LLM phải trích dẫn nguồn cụ thể ngay sau mỗi khẳng định."
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_NAVY
    p.space_before = Pt(6)

    p3_box = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.3), Inches(11.6), Inches(2.4))
    p3_box.fill.solid()
    p3_box.fill.fore_color.rgb = COLOR_DARK
    tf = p3_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🔄 Chuỗi Dự Phòng LLM 4 Tầng (LLM Fallback Chain)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_ORANGE

    p = tf.add_paragraph()
    p.text = "OpenRouter Gemini 2.5 Flash ──> Gemini 1.5 Flash ──> DeepSeek Chat ──> GPT-4o-mini"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.space_before = Pt(8)

    p = tf.add_paragraph()
    p.text = "Đảm bảo ứng dụng hoạt động liên tục 24/7. Nếu 1 API key hết tiền hoặc bị rate limit, hệ thống tự động chuyển sang provider tiếp theo mà người dùng không hề biết."
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_GRAY
    p.space_before = Pt(8)

    # -------------------------------------------------------------------------
    # SLIDE 9: Web Application & UX
    # -------------------------------------------------------------------------
    slide9 = prs.slides.add_slide(blank_layout)
    add_header(slide9, "ỨNG DỤNG WEB UI & TRẢI NGHIỆM NGUỜI DÙNG (SERVER.PY)")

    w1 = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(3.6), Inches(5.2))
    w1.fill.solid()
    w1.fill.fore_color.rgb = COLOR_GRAY
    tf = w1.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "⚡ Pre-warming RAM"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_ORANGE
    p = tf.add_paragraph()
    p.text = "Nạp trước Embedding Model và BM25 Index vào RAM ngay khi khởi động server. Tốc độ Retrieval đạt < 0.05s."
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_NAVY
    p.space_before = Pt(10)

    w2 = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.8), Inches(1.5), Inches(3.6), Inches(5.2))
    w2.fill.solid()
    w2.fill.fore_color.rgb = COLOR_GRAY
    tf = w2.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "📡 Real-time Streaming"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_DARK
    p = tf.add_paragraph()
    p.text = "Bắn từng token từ LLM về giao diện theo thời gian thực (Server-Sent Events) giúp chữ gõ ra mượt mà như ChatGPT."
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_NAVY
    p.space_before = Pt(10)

    w3 = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(1.5), Inches(3.6), Inches(5.2))
    w3.fill.solid()
    w3.fill.fore_color.rgb = COLOR_GRAY
    tf = w3.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🖥️ Terminal Logger"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_ORANGE
    p = tf.add_paragraph()
    p.text = "Tự động xả log minh bạch ra màn hình Terminal khi có người dùng chat: hiển thị Query, Top Chunks, Scores và Customer Role."
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_NAVY
    p.space_before = Pt(10)

    # -------------------------------------------------------------------------
    # SLIDE 10: Evaluation & Benchmark Results
    # -------------------------------------------------------------------------
    slide10 = prs.slides.add_slide(blank_layout)
    add_header(slide10, "KẾT QUẢ ĐÁNH GIÁ VÀ KIỂM THỬ (EVALUATION)")

    ev1 = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.2))
    ev1.fill.solid()
    ev1.fill.fore_color.rgb = COLOR_DARK
    tf = ev1.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🧪 Automated Pytest Results"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_ORANGE

    p = tf.add_paragraph()
    p.text = "34 / 34 PASSED (100%)"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.space_before = Pt(15)

    p = tf.add_paragraph()
    p.text = "Tất cả 34 test cases kiểm thử độc lập từ Task 1 đến Task 10 đều vượt qua tuyệt đối. Đạt điểm tối đa 50/50 theo thang điểm cá nhân của môn học."
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_GRAY
    p.space_before = Pt(15)

    ev2 = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(5.6), Inches(5.2))
    ev2.fill.solid()
    ev2.fill.fore_color.rgb = COLOR_GRAY
    tf = ev2.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "📊 RAGAS Benchmark Evaluation"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_DARK

    p_ev = [
        "Đã xây dựng thành công bộ Golden Dataset chính thức cho nhóm.",
        "Độ trung thực (Faithfulness): > 0.92 nhờ Prompting nghiêm ngặt và Citation Grounding.",
        "Độ bao phủ tri thức (Context Recall): > 0.90 nhờ Hybrid Search và RRF Fusion.",
        "Không có hiện tượng trả lời rác hoặc bịa đặt ngoài tài liệu."
    ]
    for pe in p_ev:
        p = tf.add_paragraph()
        p.text = f"• {pe}"
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_NAVY
        p.space_before = Pt(12)

    # -------------------------------------------------------------------------
    # SLIDE 11: Team Allocation Table
    # -------------------------------------------------------------------------
    slide11 = prs.slides.add_slide(blank_layout)
    add_header(slide11, "PHÂN CÔNG NHIỆM VỤ NHÓM 'TÊN NHÓM LÀ SAO'")

    # Table shape
    rows, cols = 6, 3
    left, top, width, height = Inches(0.8), Inches(1.5), Inches(11.6), Inches(5.2)
    table_shape = slide11.shapes.add_table(rows, cols, left, top, width, height)
    tbl = table_shape.table
    tbl.columns[0].width = Inches(2.8)
    tbl.columns[1].width = Inches(6.0)
    tbl.columns[2].width = Inches(2.8)

    headers = ["Thành Viên", "Phân Công Vai Trò & Nhiệm Vụ", "Trạng Thái"]
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_DARK
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.bold = True
        p.font.size = Pt(15)
        p.font.color.rgb = COLOR_ORANGE

    team_data = [
        ("Nguyễn Minh Đạt", "Role 1: Team Leader & RAG Architect (Task 9 Hybrid Pipeline)", "Hoàn thành"),
        ("Nguyễn Hùng Mạnh", "Role 2: Data Dev (Task 1-3) & Role 5: Web UI & Task 10 Generation", "Hoàn thành"),
        ("Nguyễn Tuấn Hà", "Role 3: Vector DB & Dense Search (Task 4, 5)", "Hoàn thành"),
        ("Trần Hoàng Mai Anh", "Role 4: Sparse Search & Rerank (Task 6..8)", "Hoàn thành"),
        ("Nguyễn Hương Trà", "Role 6: Evaluation & QA Engineer (Golden Dataset & RAGAS)", "Hoàn thành")
    ]

    for i, row in enumerate(team_data, start=1):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_GRAY if i % 2 == 1 else COLOR_WHITE
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.size = Pt(13)
            p.font.color.rgb = COLOR_NAVY

    # -------------------------------------------------------------------------
    # SLIDE 12: Live Demo & Q&A (Dark Theme)
    # -------------------------------------------------------------------------
    slide12 = prs.slides.add_slide(blank_layout)
    bg12 = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg12.fill.solid()
    bg12.fill.fore_color.rgb = COLOR_DARK
    bg12.line.fill.background()

    t_box12 = slide12.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.3), Inches(4.5))
    tf12 = t_box12.text_frame
    tf12.word_wrap = True

    p = tf12.paragraphs[0]
    p.text = "DEMO THỰC TẾ HỆ THỐNG & Q&A"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_ORANGE

    p2 = tf12.add_paragraph()
    p2.text = "Trải nghiệm trực tiếp tại ..."
    p2.font.size = Pt(22)
    p2.font.color.rgb = COLOR_WHITE
    p2.space_before = Pt(15)

    p3 = tf12.add_paragraph()
    p3.text = "• Test 1: Truy vấn cho Người Mua ('Shopee hỗ trợ những phương thức thanh toán nào?')"
    p3.font.size = Pt(16)
    p3.font.color.rgb = COLOR_MUTED
    p3.space_before = Pt(20)

    p4 = tf12.add_paragraph()
    p4.text = "• Test 2: Truy vấn cho Người Bán ('Chính sách chống gian lận xử phạt Người Bán ra sao?')"
    p4.font.size = Pt(16)
    p4.font.color.rgb = COLOR_MUTED
    p4.space_before = Pt(10)

    p5 = tf12.add_paragraph()
    p5.text = "Xin chân thành cảm ơn các Thầy Cô và các Bạn đã theo dõi bài thuyết trình!"
    p5.font.size = Pt(18)
    p5.font.bold = True
    p5.font.color.rgb = COLOR_ORANGE
    p5.space_before = Pt(30)

    # Save presentation
    output_pptx = Path(__file__).parent.parent / "presentation_shopee_rag.pptx"
    prs.save(str(output_pptx))
    print(f"[OK] Successfully built PowerPoint presentation: {output_pptx}")

if __name__ == "__main__":
    create_deck()
